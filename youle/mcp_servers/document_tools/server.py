"""mcp-document-tools:pptx / xlsx / docx / pdf / pdf_ocr。

设计原则:
- **优雅降级**:可选依赖未安装时返回 `_mock=True`,不让 server 启动失败
- **真实实现**:python-pptx / openpyxl / python-docx / pdfplumber 接入
- **OSS 落盘**:产物全部写 OSS,返回 `oss_ref`(铁律 4)
- **铁律 16 / V1.5 留接口**:复杂 PPT 排版能力归 V1.5;此处先做 happy path
"""

from __future__ import annotations

import asyncio
import io
import os
import tempfile
from pathlib import Path
from typing import Any
from urllib.parse import urlparse
from uuid import uuid4

import boto3
import httpx
import structlog
from botocore.client import Config

from mcp_servers._shared.http_app import make_app

log = structlog.get_logger(__name__)

OSS_ENDPOINT = os.getenv("OSS_ENDPOINT", "http://minio:9000")
OSS_BUCKET = os.getenv("OSS_BUCKET", "youle-dev")


def _s3():
    return boto3.client(
        "s3",
        endpoint_url=OSS_ENDPOINT,
        aws_access_key_id=os.getenv("OSS_ACCESS_KEY", "minioadmin"),
        aws_secret_access_key=os.getenv("OSS_SECRET_KEY", "minioadmin"),
        config=Config(signature_version="s3v4"),
    )


def _read_ref(ref: str) -> bytes:
    if ref.startswith("oss://"):
        parsed = urlparse(ref)
        bucket = parsed.netloc
        key = parsed.path.lstrip("/")
        obj = _s3().get_object(Bucket=bucket, Key=key)
        return obj["Body"].read()
    if ref.startswith(("http://", "https://")):
        with httpx.Client(timeout=30.0, follow_redirects=True) as client:
            return client.get(ref).content
    return Path(ref).read_bytes()


def _put_oss(key: str, data: bytes, content_type: str) -> str:
    _s3().put_object(Bucket=OSS_BUCKET, Key=key, Body=data, ContentType=content_type)
    return f"oss://{OSS_BUCKET}/{key}"


# ─────────────────────── PPTX ───────────────────────
def _do_pptx_assemble(
    *,
    title: str | None,
    slides: list[dict[str, Any]],
) -> tuple[str, int]:
    """slides: [{title, bullets, image_ref?}]"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        raise RuntimeError("python-pptx not installed")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 封面
    if title:
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = title

    body_layout = prs.slide_layouts[1]
    for sd in slides:
        slide = prs.slides.add_slide(body_layout)
        if slide.shapes.title and sd.get("title"):
            slide.shapes.title.text = str(sd["title"])
        bullets = sd.get("bullets") or []
        # 找内容占位符
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                tf = shape.text_frame
                tf.clear()
                for i, b in enumerate(bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = str(b)
                    p.font.size = Pt(20)
                break
        # 可选图片
        img = sd.get("image_ref")
        if img:
            try:
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                    tmp.write(_read_ref(img))
                    tmp_path = tmp.name
                slide.shapes.add_picture(
                    tmp_path, Inches(8), Inches(2), width=Inches(4.5)
                )
                Path(tmp_path).unlink(missing_ok=True)
            except Exception as e:
                log.warning("pptx.image_failed", err=str(e))

    buf = io.BytesIO()
    prs.save(buf)
    data = buf.getvalue()
    key = f"docs/pptx/{uuid4().hex}.pptx"
    return _put_oss(
        key,
        data,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ), len(data)


async def pptx_assemble(arguments: dict[str, Any]) -> dict[str, Any]:
    title = arguments.get("title")
    slides = arguments.get("slides") or []
    if not slides:
        return {"error": "slides 必填", "_failed": True}
    try:
        oss_ref, size = await asyncio.to_thread(
            _do_pptx_assemble, title=title, slides=slides
        )
        return {"oss_ref": oss_ref, "slide_count": len(slides), "size_bytes": size}
    except RuntimeError as e:
        log.warning("pptx.lib_missing", err=str(e))
        return {"oss_ref": f"oss://{OSS_BUCKET}/mock-output.pptx", "_mock": True}
    except Exception as e:
        log.exception("pptx.failed", err=str(e))
        return {"error": str(e), "_failed": True}


# ─────────────────────── XLSX ───────────────────────
def _do_xlsx_assemble(*, sheets: list[dict[str, Any]]) -> tuple[str, int]:
    try:
        from openpyxl import Workbook
    except ImportError:
        raise RuntimeError("openpyxl not installed")

    wb = Workbook()
    wb.remove(wb.active)
    for sd in sheets:
        ws = wb.create_sheet(title=str(sd.get("name", "Sheet1"))[:31])
        rows: list[list[Any]] = sd.get("rows") or []
        headers = sd.get("headers")
        if headers:
            ws.append([str(h) for h in headers])
        for row in rows:
            ws.append(list(row))

    buf = io.BytesIO()
    wb.save(buf)
    data = buf.getvalue()
    key = f"docs/xlsx/{uuid4().hex}.xlsx"
    return _put_oss(
        key,
        data,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ), len(data)


async def xlsx_assemble(arguments: dict[str, Any]) -> dict[str, Any]:
    sheets = arguments.get("sheets") or []
    # 兼容 简单调用:rows + headers
    if not sheets and arguments.get("rows"):
        sheets = [
            {
                "name": arguments.get("sheet_name", "Sheet1"),
                "headers": arguments.get("headers"),
                "rows": arguments.get("rows"),
            }
        ]
    if not sheets:
        return {"error": "sheets 必填", "_failed": True}
    try:
        oss_ref, size = await asyncio.to_thread(_do_xlsx_assemble, sheets=sheets)
        return {"oss_ref": oss_ref, "sheet_count": len(sheets), "size_bytes": size}
    except RuntimeError as e:
        log.warning("xlsx.lib_missing", err=str(e))
        return {"oss_ref": f"oss://{OSS_BUCKET}/mock-output.xlsx", "_mock": True}
    except Exception as e:
        log.exception("xlsx.failed", err=str(e))
        return {"error": str(e), "_failed": True}


# ─────────────────────── DOCX ───────────────────────
def _do_docx_assemble(*, title: str | None, paragraphs: list[Any]) -> tuple[str, int]:
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx not installed")

    doc = Document()
    if title:
        doc.add_heading(str(title), level=0)
    for p in paragraphs:
        if isinstance(p, dict):
            kind = p.get("kind", "p")
            text = str(p.get("text", ""))
            if kind in ("h1", "h2", "h3"):
                doc.add_heading(text, level=int(kind[1]))
            elif kind == "list":
                doc.add_paragraph(text, style="List Bullet")
            else:
                doc.add_paragraph(text)
        else:
            doc.add_paragraph(str(p))

    buf = io.BytesIO()
    doc.save(buf)
    data = buf.getvalue()
    key = f"docs/docx/{uuid4().hex}.docx"
    return _put_oss(
        key, data, "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ), len(data)


async def docx_assemble(arguments: dict[str, Any]) -> dict[str, Any]:
    title = arguments.get("title")
    paragraphs = arguments.get("paragraphs") or []
    if not paragraphs and arguments.get("text"):
        paragraphs = [str(arguments["text"])]
    if not paragraphs:
        return {"error": "paragraphs 必填", "_failed": True}
    try:
        oss_ref, size = await asyncio.to_thread(
            _do_docx_assemble, title=title, paragraphs=paragraphs
        )
        return {"oss_ref": oss_ref, "paragraph_count": len(paragraphs), "size_bytes": size}
    except RuntimeError as e:
        log.warning("docx.lib_missing", err=str(e))
        return {"oss_ref": f"oss://{OSS_BUCKET}/mock-output.docx", "_mock": True}
    except Exception as e:
        log.exception("docx.failed", err=str(e))
        return {"error": str(e), "_failed": True}


# ─────────────────────── PDF extract ───────────────────────
def _do_pdf_extract(ref: str, max_pages: int) -> dict[str, Any]:
    try:
        import pdfplumber
    except ImportError:
        raise RuntimeError("pdfplumber not installed")

    raw = _read_ref(ref)
    pages_text: list[str] = []
    with pdfplumber.open(io.BytesIO(raw)) as pdf:
        total = len(pdf.pages)
        for i, page in enumerate(pdf.pages):
            if i >= max_pages:
                break
            pages_text.append(page.extract_text() or "")
    return {"pages": pages_text, "page_count_total": total}


async def pdf_extract(arguments: dict[str, Any]) -> dict[str, Any]:
    ref = arguments.get("ref") or arguments.get("oss_ref") or arguments.get("url")
    if not ref:
        return {"error": "ref 必填", "_failed": True}
    max_pages = int(arguments.get("max_pages", 50))
    try:
        result = await asyncio.to_thread(_do_pdf_extract, ref, max_pages)
        text = "\n\n".join(result["pages"])
        return {
            "text": text,
            "char_count": len(text),
            "page_count": len(result["pages"]),
            "page_count_total": result["page_count_total"],
        }
    except RuntimeError as e:
        log.warning("pdf_extract.lib_missing", err=str(e))
        return {"text": "[mock] extracted text", "_mock": True}
    except Exception as e:
        log.exception("pdf_extract.failed", err=str(e))
        return {"error": str(e), "_failed": True}


# ─────────────────────── PDF OCR ───────────────────────
def _do_pdf_ocr(ref: str, max_pages: int, lang: str) -> dict[str, Any]:
    """优先 pytesseract,其次返回 _mock。V1.5 接阿里云 OCR。"""
    try:
        import pytesseract
        from pdf2image import convert_from_bytes
    except ImportError:
        raise RuntimeError("pytesseract or pdf2image not installed")

    raw = _read_ref(ref)
    images = convert_from_bytes(raw)
    pages_text: list[str] = []
    for i, img in enumerate(images):
        if i >= max_pages:
            break
        pages_text.append(pytesseract.image_to_string(img, lang=lang))
    return {"pages": pages_text}


async def pdf_ocr(arguments: dict[str, Any]) -> dict[str, Any]:
    ref = arguments.get("ref") or arguments.get("oss_ref") or arguments.get("url")
    if not ref:
        return {"error": "ref 必填", "_failed": True}
    max_pages = int(arguments.get("max_pages", 20))
    lang = str(arguments.get("lang", "chi_sim+eng"))
    try:
        result = await asyncio.to_thread(_do_pdf_ocr, ref, max_pages, lang)
        text = "\n\n".join(result["pages"])
        return {
            "text": text,
            "char_count": len(text),
            "page_count": len(result["pages"]),
        }
    except RuntimeError as e:
        log.warning("pdf_ocr.lib_missing", err=str(e))
        return {"text": "[mock] ocr text", "_mock": True}
    except Exception as e:
        log.exception("pdf_ocr.failed", err=str(e))
        return {"error": str(e), "_failed": True}


app = make_app(
    server_name="document-tools",
    tools={
        "pptx_assemble": pptx_assemble,
        "xlsx_assemble": xlsx_assemble,
        "docx_assemble": docx_assemble,
        "pdf_extract": pdf_extract,
        "pdf_ocr": pdf_ocr,
    },
)


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=7005)
